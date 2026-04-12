from pathlib import Path
from typing import Dict, Any, List
import json

import pandas as pd
from docx import Document
from pypdf import PdfReader
from pptx import Presentation

from app.layers.file_ingestion.file_chunker import split_text_into_chunks
from app.layers.file_ingestion.ingestion_schema import ParsedFile


class FileParser:
    def parse(self, file_path: str, file_id: str) -> ParsedFile:
        path = Path(file_path)
        suffix = path.suffix.lower()

        if suffix == ".pdf":
            text = self._parse_pdf(path)
            chunks = split_text_into_chunks(text)
            return ParsedFile(
                file_id=file_id,
                file_name=path.name,
                file_type=suffix,
                raw_text=text,
                chunks=chunks,
                structured_summary=None,
            )

        if suffix in [".txt", ".md"]:
            text = self._parse_text(path)
            chunks = split_text_into_chunks(text)
            return ParsedFile(
                file_id=file_id,
                file_name=path.name,
                file_type=suffix,
                raw_text=text,
                chunks=chunks,
                structured_summary=None,
            )

        if suffix == ".docx":
            text = self._parse_docx(path)
            chunks = split_text_into_chunks(text)
            return ParsedFile(
                file_id=file_id,
                file_name=path.name,
                file_type=suffix,
                raw_text=text,
                chunks=chunks,
                structured_summary=None,
            )

        if suffix == ".pptx":
            text = self._parse_pptx(path)
            chunks = split_text_into_chunks(text)
            return ParsedFile(
                file_id=file_id,
                file_name=path.name,
                file_type=suffix,
                raw_text=text,
                chunks=chunks,
                structured_summary=None,
            )

        if suffix in [".csv", ".xlsx", ".xls"]:
            df = self._parse_table(path, suffix)
            summary = self._summarize_dataframe(df, path.name)
            summary_text = json.dumps(summary, ensure_ascii=False, indent=2)
            chunks = split_text_into_chunks(summary_text, chunk_size=1800, overlap=200)

            return ParsedFile(
                file_id=file_id,
                file_name=path.name,
                file_type=suffix,
                raw_text=summary_text,
                chunks=chunks,
                structured_summary=summary,
            )

        raise ValueError(f"Unsupported file type: {suffix}")

    def _parse_text(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="ignore")

    def _parse_docx(self, path: Path) -> str:
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)

    def _parse_pdf(self, path: Path) -> str:
        reader = PdfReader(str(path))
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text() or "")
        return "\n".join(texts)

    def _parse_pptx(self, path: Path) -> str:
        prs = Presentation(str(path))
        slides_text = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            parts = [f"--- Slide {slide_idx} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip():
                            parts.append(row_text)
            slides_text.append("\n".join(parts))
        return "\n\n".join(slides_text)

    def _parse_table(self, path: Path, suffix: str) -> pd.DataFrame:
        if suffix == ".csv":
            return pd.read_csv(path)

        if suffix == ".xlsx":
            return pd.read_excel(path, engine="openpyxl")

        if suffix == ".xls":
            return pd.read_excel(path, engine="xlrd")

        raise ValueError(f"Unsupported table file type: {suffix}")

    def _summarize_dataframe(self, df: pd.DataFrame, file_name: str) -> Dict[str, Any]:
        df = df.copy()

        num_rows, num_cols = df.shape
        columns = [str(c) for c in df.columns.tolist()]

        dtypes = {str(col): str(dtype) for col, dtype in df.dtypes.items()}
        missing_summary = {str(col): int(df[col].isna().sum()) for col in df.columns}

        sample_rows = df.head(5).fillna("").to_dict(orient="records")

        numeric_columns = df.select_dtypes(include=["number"]).columns.tolist()
        categorical_columns = [c for c in df.columns if c not in numeric_columns]

        numeric_summary = {}
        for col in numeric_columns[:20]:
            series = df[col].dropna()
            if len(series) == 0:
                numeric_summary[str(col)] = {
                    "count": 0,
                    "mean": None,
                    "std": None,
                    "min": None,
                    "max": None,
                }
            else:
                numeric_summary[str(col)] = {
                    "count": int(series.count()),
                    "mean": float(series.mean()) if pd.notna(series.mean()) else None,
                    "std": float(series.std()) if pd.notna(series.std()) else None,
                    "min": float(series.min()) if pd.notna(series.min()) else None,
                    "max": float(series.max()) if pd.notna(series.max()) else None,
                }

        categorical_summary = {}
        for col in categorical_columns[:20]:
            series = df[col].dropna().astype(str)
            top_values = series.value_counts().head(10)
            categorical_summary[str(col)] = {
                "unique_count": int(series.nunique()),
                "top_values": top_values.to_dict(),
            }

        possible_targets = self._guess_possible_targets(df)

        return {
            "file_name": file_name,
            "dataset_overview": {
                "num_rows": int(num_rows),
                "num_cols": int(num_cols),
                "columns": columns,
                "dtypes": dtypes,
            },
            "missing_summary": missing_summary,
            "sample_rows": sample_rows,
            "numeric_summary": numeric_summary,
            "categorical_summary": categorical_summary,
            "possible_target_columns": possible_targets,
        }

    def _guess_possible_targets(self, df: pd.DataFrame) -> List[str]:
        candidates = []

        for col in df.columns:
            name = str(col).lower()

            if any(keyword in name for keyword in [
                "target", "label", "class", "outcome", "y", "price", "sale", "sold"
            ]):
                candidates.append(str(col))
                continue

            non_null = df[col].dropna()
            if len(non_null) == 0:
                continue

            unique_ratio = non_null.nunique() / max(len(non_null), 1)

            # 类别型目标常见：唯一值不太多
            if non_null.nunique() <= 10:
                candidates.append(str(col))
                continue

            # 回归目标：数值列且名字像结果列
            if pd.api.types.is_numeric_dtype(df[col]) and unique_ratio > 0.2:
                if any(keyword in name for keyword in [
                    "score", "amount", "value", "price", "income"
                ]):
                    candidates.append(str(col))

        # 去重
        deduped = []
        for c in candidates:
            if c not in deduped:
                deduped.append(c)

        return deduped[:10]